from reportlab.pdfgen import canvas
from pptx import Presentation

class Writer:
    def create_outputs(self, literature, plan, analysis, session_id):

        # ---------- CREATE PDF ----------
        pdf_path = f"artifacts/{session_id}/paper.pdf"
        c = canvas.Canvas(pdf_path)

        c.setFont("Helvetica-Bold", 20)
        c.drawString(50, 800, "Research Paper")
        
        c.setFont("Helvetica", 12)
        c.drawString(50, 770, f"Topic: {plan['goal']}")
        c.drawString(50, 750, f"Model Used: {plan['model']}")
        c.drawString(50, 730, f"Metric: {plan['metric']}")
        c.drawString(50, 710, f"Accuracy: {analysis['accuracy']}")
        
        c.drawString(50, 680, "Literature Summary:")
        y = 660
        for p in literature:
            c.drawString(60, y, f"- {p['title']}: {p['summary']}")
            y -= 20
        
        c.save()


        # ---------- CREATE PPT ----------
        ppt_path = f"artifacts/{session_id}/slides.pptx"
        prs = Presentation()

        # Title Slide
        slide_1 = prs.slides.add_slide(prs.slide_layouts[0])
        slide_1.shapes.title.text = "ResearchLab-AI Results"
        slide_1.placeholders[1].text = "Automatically generated research presentation"

        # Slide 2 - Experiment
        slide_2 = prs.slides.add_slide(prs.slide_layouts[1])
        slide_2.shapes.title.text = "Experiment Overview"
        slide_2.placeholders[1].text = f"Model: {plan['model']}\nMetric: {plan['metric']}\nAccuracy: {analysis['accuracy']}"

        # Slide 3 - Literature
        slide_3 = prs.slides.add_slide(prs.slide_layouts[1])
        slide_3.shapes.title.text = "Literature Summary"
        slide_3.placeholders[1].text = "\n".join([f"- {p['title']}" for p in literature])

        prs.save(ppt_path)

        return {"paper": pdf_path, "slides": ppt_path}

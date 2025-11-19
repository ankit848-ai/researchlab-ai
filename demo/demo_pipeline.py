import os
import json
from sklearn.datasets import load_iris
from sklearn.model_selection import train_test_split
from sklearn.tree import DecisionTreeClassifier
from sklearn.metrics import accuracy_score
import matplotlib.pyplot as plt
from pptx import Presentation
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas

# ---------------------------
# Create artifacts folder
# ---------------------------
ARTIFACT_DIR = "../artifacts/session_001/"
os.makedirs(ARTIFACT_DIR, exist_ok=True)

# ---------------------------
# Step 1 — Load dataset
# ---------------------------
iris = load_iris()
X = iris.data
y = iris.target

X_train, X_test, y_train, y_test = train_test_split(
    X, y, test_size=0.2, random_state=42
)

# ---------------------------
# Step 2 — Train model
# ---------------------------
clf = DecisionTreeClassifier()
clf.fit(X_train, y_train)

# ---------------------------
# Step 3 — Predictions & Metrics
# ---------------------------
preds = clf.predict(X_test)
acc = accuracy_score(y_test, preds)

metrics = {"accuracy": acc}
json.dump(metrics, open(ARTIFACT_DIR + "metrics.json", "w"))

# ---------------------------
# Step 4 — Plot accuracy
# ---------------------------
plt.figure(figsize=(6, 4))
plt.bar(["DecisionTree Accuracy"], [acc])
plt.ylim(0, 1)
plt.ylabel("Accuracy")
plt.title("Model Accuracy")
plt.savefig(ARTIFACT_DIR + "accuracy.png")
plt.close()

# ---------------------------
# Step 5 — Create PDF (Research Paper Mock)
# ---------------------------
pdf_path = ARTIFACT_DIR + "paper.pdf"
c = canvas.Canvas(pdf_path, pagesize=letter)
c.drawString(100, 750, "ResearchLab-AI: Auto-generated Research Paper")
c.drawString(100, 720, f"Model Used: DecisionTreeClassifier")
c.drawString(100, 700, f"Accuracy: {acc}")
c.save()

# ---------------------------
# Step 6 — Create PPTX Slide Deck
# ---------------------------
ppt = Presentation()
slide = ppt.slides.add_slide(ppt.slide_layouts[1])
title = slide.shapes.title
body = slide.placeholders[1]

title.text = "ResearchLab-AI Presentation"
body.text = f"Model Accuracy: {acc}"

ppt.save(ARTIFACT_DIR + "slides.pptx")

# ---------------------------
# Step 7 — Log file
# ---------------------------
with open(ARTIFACT_DIR + "run_logs.txt", "w") as f:
    f.write("Pipeline executed successfully.\n")
    f.write(f"Accuracy: {acc}\n")

print("Artifacts generated successfully at", ARTIFACT_DIR)
